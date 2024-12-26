import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from anthropic import Anthropic
import PyPDF2
import io

# Initialize Anthropic client
client = Anthropic(api_key=st.secrets["ANTHROPIC_API_KEY"])

def extract_text_from_pdf(pdf_file):
    """Extract text from uploaded PDF"""
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def create_mckinsey_slide(prs, title, content):
    """Create a slide following McKinsey style"""
    # McKinsey style slide layout (Title and Content)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Title style
    title_shape = slide.shapes.title
    title_shape.text = title
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.font.name = 'Helvetica'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Content style
    content_shape = slide.shapes.placeholders[1]
    content_frame = content_shape.text_frame
    p = content_frame.add_paragraph()
    p.text = content
    p.font.name = 'Helvetica'
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    return slide

def generate_insights(text: str) -> dict:
    """Generate structured insights using Claude"""
    prompt = f"""
    You are a McKinsey consultant. Based on the following text, create a structured presentation outline:

    Text:
    {text}

    Please provide:
    1. An executive summary (2-3 key points)
    2. Main findings (3-4 points with supporting details)
    3. Recommendations (2-3 actionable items)
    4. Next steps (2-3 concrete actions)

    Format each section as clear, concise bullet points suitable for a presentation.
    """
    
    response = client.messages.create(
        model="claude-3-sonnet-20240229",
        max_tokens=1500,
        messages=[{
            "role": "user",
            "content": prompt
        }]
    )
    
    return response.content[0].text

def create_presentation(insights):
    """Create McKinsey-style presentation"""
    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Strategic Analysis"
    subtitle = title_slide.placeholders[1]
    subtitle.text = "Executive Summary"
    
    # Create content slides
    sections = insights.split('\n\n')
    for section in sections:
        if section.strip():
            lines = section.strip().split('\n')
            title = lines[0]
            content = '\n'.join(lines[1:])
            create_mckinsey_slide(prs, title, content)
    
    # Save presentation
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()

def main():
    st.set_page_config(page_title="McKinsey-Style Presentation Generator")
    
    st.title("McKinsey-Style Presentation Generator")
    st.write("Upload PDFs and generate professional presentations with AI insights")
    
    uploaded_files = st.file_uploader(
        "Upload PDF documents", 
        type=['pdf'],
        accept_multiple_files=True
    )
    
    if uploaded_files:
        combined_text = ""
        for pdf_file in uploaded_files:
            text = extract_text_from_pdf(pdf_file)
            combined_text += text + "\n\n"
        
        st.success(f"{len(uploaded_files)} file(s) uploaded successfully!")
        
        if st.button("Generate Presentation"):
            with st.spinner("Analyzing documents and generating presentation..."):
                try:
                    # Generate insights
                    insights = generate_insights(combined_text)
                    
                    # Create presentation
                    pptx_data = create_presentation(insights)
                    
                    # Show insights preview
                    st.subheader("Generated Insights")
                    st.write(insights)
                    
                    # Download button
                    st.download_button(
                        label="Download Presentation",
                        data=pptx_data,
                        file_name="mckinsey_style_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                    
                except Exception as e:
                    st.error(f"Error generating presentation: {str(e)}")

if __name__ == "__main__":
    main()
